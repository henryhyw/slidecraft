"""Inspect and edit the actual native PowerPoint in a shared component set.

The browser receives object geometry and chart/table data, never a screenshot
masquerading as an editable component. Office remains the exact rendering engine.
"""
from __future__ import annotations

import base64
import copy
import math
import shutil
import tempfile
from pathlib import Path

from . import library_sets
from .storage import ConflictError, locked, revision, update


def source(set_id: str, component_id: str):
    record = library_sets.set_record(set_id)
    if record["kind"] != "components" or record["source"] != "local":
        raise ValueError("Choose a local component set")
    catalog_path, catalog = library_sets.set_catalog(set_id)
    item = catalog["items"].get(component_id)
    if item is None:
        raise FileNotFoundError(component_id)
    relative = item.get("path") or catalog.get("native_source")
    path = (catalog_path.parent / relative).resolve() if relative else None
    if path and (catalog_path.parent.resolve() not in path.parents or not path.is_file()):
        raise FileNotFoundError("Component source is missing or outside its set")
    return catalog_path, item, path


def color(fill, fallback=None):
    try:
        return "#" + str(fill.fore_color.rgb) if fill.fore_color.type == 1 else fallback
    except (AttributeError, TypeError, ValueError):
        return fallback


def font_data(font):
    try:
        rgb = "#" + str(font.color.rgb) if font.color.type == 1 else None
    except (AttributeError, TypeError, ValueError):
        rgb = None
    return {"font": font.name, "font_size": font.size.pt if font.size else None, "bold": font.bold, "italic": font.italic, "color": rgb}


def shape_data(shape):
    item = {"id": str(shape.shape_id), "name": shape.name, "x": shape.left.pt, "y": shape.top.pt,
            "width": shape.width.pt, "height": shape.height.pt, "rotation": shape.rotation,
            "type": "shape", "fill": color(getattr(shape, "fill", None)),
            "line": color(getattr(getattr(shape, "line", None), "fill", None))}
    if shape.shape_type == 6:
        item.update(type="group", children=[shape_data(child) for child in shape.shapes])
        # Child coordinates are in the group's local drawing space.
        xfrm = shape.element.grpSpPr.xfrm
        if xfrm is not None and xfrm.chOff is not None and xfrm.chExt is not None:
            item["child_space"] = [xfrm.chOff.x / 12700, xfrm.chOff.y / 12700, xfrm.chExt.cx / 12700, xfrm.chExt.cy / 12700]
    elif shape.has_chart:
        chart = shape.chart
        chart_type = chart.chart_type.name
        categories = [category.label for category in chart.plots[0].categories] if hasattr(chart.plots[0], "categories") else []
        series = []
        for sequence in chart.series:
            series.append({"name": sequence.name, "values": list(sequence.values), "color": color(sequence.format.fill),
                           "point_colors": [color(point.format.fill) for point in sequence.points]})
        item.update(type="chart", chart={"type": chart_type, "categories": categories, "series": series,
                    "show_legend": chart.has_legend, "title": chart.chart_title.text_frame.text if chart.has_title and chart.chart_title.has_text_frame else ""})
        item["chart"]["editable_data"] = not any(key in chart_type for key in ("XY_", "BUBBLE", "STOCK")) and len(chart.plots) == 1
        for key in ("hole_size", "gap_width", "overlap"):
            if hasattr(chart.plots[0], key):
                item["chart"][key] = getattr(chart.plots[0], key)
    elif shape.has_table:
        table = shape.table
        item.update(type="table", cells=[[{"text": cell.text, "fill": color(cell.fill),
                            **font_data(cell.text_frame.paragraphs[0].font),
                            "merged": cell.is_spanned, "row_span": cell.span_height, "col_span": cell.span_width}
                            for cell in row.cells] for row in table.rows],
                    column_widths=[column.width.pt for column in table.columns], row_heights=[row.height.pt for row in table.rows])
    elif shape.shape_type == 13:
        item.update(type="image", image=f"data:{shape.image.content_type};base64,{base64.b64encode(shape.image.blob).decode()}")
    elif shape.has_text_frame:
        item.update(type="text" if shape.text.strip() else "shape", text=shape.text,
                    paragraphs=[{"text": paragraph.text, **font_data(paragraph.font),
                        "runs": [{"text": run.text, **font_data(run.font)} for run in paragraph.runs]}
                        for paragraph in shape.text_frame.paragraphs])
    try:
        item["geometry"] = shape.auto_shape_type.name
    except (ValueError, AttributeError):
        pass
    return item


def inspect(set_id: str, component_id: str, slide_number: int | None = None):
    catalog_path, item, path = source(set_id, component_id)
    result = {"set_id": set_id, "component": copy.deepcopy(item), "catalog_revision": revision(catalog_path)}
    if path is None or path.suffix.lower() != ".pptx":
        return {**result, "native": False, "note": "This entry defines a design grammar. It has no native PowerPoint source yet."}
    from pptx import Presentation
    deck = Presentation(path)
    number = slide_number or int(item.get("native_source_slide_number", 1))
    if not 1 <= number <= len(deck.slides):
        raise ValueError("Slide number is outside the component source")
    slide = deck.slides[number - 1]
    return {**result, "native": True, "source_revision": revision(path), "source_name": path.name,
            "slide_number": number, "slide_count": len(deck.slides), "width": deck.slide_width.pt,
            "height": deck.slide_height.pt, "background": color(slide.background.fill, "#FFFFFF"),
            "objects": [shape_data(shape) for shape in slide.shapes]}


def update_definition(set_id, component_id, values, expected):
    path, item, _ = source(set_id, component_id)
    if not isinstance(values, dict) or {"id", "path", "preview_path", "native_source_slide_number"} & set(values):
        raise ValueError("Edit component guidance without changing source identity")
    def change(catalog):
        catalog["items"][component_id].update(values)
        return catalog
    update(path, change, expected=expected)
    return inspect(set_id, component_id)


def save_object(set_id, component_id, object_id, values, expected, slide_number=None):
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    _, item, path = source(set_id, component_id)
    if path is None or path.suffix.lower() != ".pptx":
        raise ValueError("This component has no editable PowerPoint source")
    allowed = {"name", "x", "y", "width", "height", "rotation", "text", "fill", "line", "font", "font_size", "color", "bold", "italic", "chart", "cells"}
    if set(values) - allowed:
        raise ValueError("Unsupported object property")
    with locked(path):
        if revision(path) != expected:
            raise ConflictError("This PowerPoint changed. Reopen the component before saving.")
        deck = Presentation(path)
        number = slide_number or int(item.get("native_source_slide_number", 1))
        if not 1 <= number <= len(deck.slides):
            raise ValueError("Invalid source slide")
        def flatten(shapes):
            for shape in shapes:
                yield shape
                if shape.shape_type == 6:
                    yield from flatten(shape.shapes)
        shape = next((shape for shape in flatten(deck.slides[number - 1].shapes) if str(shape.shape_id) == str(object_id)), None)
        if shape is None:
            raise ValueError("Object no longer exists")
        for key, prop in {"x": "left", "y": "top", "width": "width", "height": "height"}.items():
            if key in values:
                number_value = float(values[key])
                if not math.isfinite(number_value) or (key in {"width", "height"} and number_value <= 0):
                    raise ValueError("Use finite geometry with positive dimensions")
                setattr(shape, prop, Pt(number_value))
        for key in ("name", "rotation"):
            if key in values:
                setattr(shape, key, values[key])
        for key in ("fill", "line"):
            if key in values and values[key]:
                fill = shape.fill if key == "fill" else shape.line.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(values[key].lstrip("#"))
        if "text" in values:
            if not shape.has_text_frame:
                raise ValueError("This object has no text frame")
            if shape.text != values["text"]:
                # Preserve styling of corresponding paragraphs and runs where possible.
                old = shape.text_frame.paragraphs[0].font
                style = font_data(old)
                shape.text = values["text"]
                for paragraph in shape.text_frame.paragraphs:
                    if style["font"]:
                        paragraph.font.name = style["font"]
                    if style["font_size"]:
                        paragraph.font.size = Pt(style["font_size"])
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for font in [paragraph.font, *(run.font for run in paragraph.runs)]:
                    for key, attr in (("font", "name"), ("bold", "bold"), ("italic", "italic")):
                        if key in values:
                            setattr(font, attr, values[key])
                    if values.get("font_size"):
                        font.size = Pt(float(values["font_size"]))
                    if values.get("color"):
                        font.color.rgb = RGBColor.from_string(values["color"].lstrip("#"))
        if "cells" in values:
            cells = values["cells"]
            if not shape.has_table or len(cells) != len(shape.table.rows) or any(len(row) != len(shape.table.columns) for row in cells):
                raise ValueError("Keep the source table dimensions when editing cell contents")
            for r, row in enumerate(cells):
                for c, cell in enumerate(row):
                    target = shape.table.cell(r, c)
                    if not target.is_spanned and target.text != cell["text"]:
                        target.text = cell["text"]
        if "chart" in values:
            if not shape.has_chart or not shape_data(shape)["chart"]["editable_data"]:
                raise ValueError("Edit this chart type directly in PowerPoint")
            chart = values["chart"]
            data = CategoryChartData()
            data.categories = chart["categories"]
            for sequence in chart["series"]:
                numbers = sequence["values"]
                if len(numbers) != len(chart["categories"]) or any(v is not None and (isinstance(v, bool) or not isinstance(v, (float, int)) or not math.isfinite(v)) for v in numbers):
                    raise ValueError("Every series needs one finite number or empty value per category")
                data.add_series(sequence["name"], numbers)
            shape.chart.replace_data(data)
            shape.chart.has_legend = chart.get("show_legend", shape.chart.has_legend)
            if "title" in chart:
                shape.chart.has_title = bool(chart["title"])
                if chart["title"]:
                    shape.chart.chart_title.text_frame.text = chart["title"]
            for key in ("hole_size", "gap_width", "overlap"):
                if key in chart and hasattr(shape.chart.plots[0], key):
                    setattr(shape.chart.plots[0], key, int(chart[key]))
            for series, entry in zip(shape.chart.series, chart["series"]):
                if entry.get("color"):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor.from_string(entry["color"].lstrip("#"))
        backup = path.parent / ".history" / f"{path.stem}-{expected[:16]}.pptx"
        backup.parent.mkdir(exist_ok=True)
        if not backup.exists():
            shutil.copy2(path, backup)
        with tempfile.NamedTemporaryFile(suffix=".pptx", dir=path.parent, delete=False) as file:
            temporary = Path(file.name)
        try:
            deck.save(temporary)
            temporary.replace(path)
        finally:
            temporary.unlink(missing_ok=True)
    return inspect(set_id, component_id, slide_number)
